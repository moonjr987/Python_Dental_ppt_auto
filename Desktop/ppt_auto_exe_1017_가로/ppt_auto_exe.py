import tkinter as tk
from tkinterdnd2 import DND_FILES, TkinterDnD
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.dml.color import RGBColor
import tempfile
from tkinter import filedialog
from tkinter import *
import subprocess
import sys


def resource_path(relative_path):
    base_path = getattr(sys, '_MEIPASS', os.path.dirname(os.path.abspath(__file__)))
    return os.path.join(base_path, relative_path)

# 재귀적으로 이미지 파일을 수집하는 함수
def collect_images(folder_path):
    image_files = []
    for root, dirs, files in os.walk(folder_path):
        for file in files:
            if file.lower().endswith(('.jpg', '.jpeg', '.png', '.gif')):
                image_files.append(os.path.join(root, file))
    return image_files

def open_folder_dialog():
    # image_folder = filedialog.askdirectory()  # 드래그 앤 드롭한 폴더 경로
    image_folder = filedialog.askdirectory()
    text.delete(1.0, tk.END)  # 텍스트 창 초기화
    

    # 폴더 이름(날짜) 추출
    folder_name = os.path.basename(image_folder)

    # 텍스트 상자에 폴더 이름 설정
    # text.insert(tk.END, f"이미지 폴더 경로: {image_folder}\n")
    # text.insert(tk.END, f"폴더 이름(날짜): {folder_name}\n")

    # PowerPoint 프레젠테이션 생성
    ppt = Presentation()

    

    # 슬라이드 배경색을 검은색으로 설정
    for slide in ppt.slides:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)  # 검은색 배경

    # 모든 이미지 파일 수집 (재귀적으로 하위 폴더 검색)
    subfolders = [f.path for f in os.scandir(image_folder) if f.is_dir()]

    for subfolder in subfolders:
        # 이미지 파일을 수집하고 파일 이름 순으로 정렬
        image_files = collect_images(subfolder)
        image_files.sort(key=lambda x: os.path.basename(x))

        # 슬라이드 레이아웃 생성
        layout = ppt.slide_layouts[5]  # 5는 이미지 슬라이드 레이아웃입니다.
        

        # 이미지를 그룹 단위로 슬라이드에 추가
        grouped_images = {}
        for image_file in image_files:
            filename = os.path.basename(image_file)
            group_identifier = os.path.splitext(filename)[0].split("-")[0]
            if group_identifier not in grouped_images:
                grouped_images[group_identifier] = []
            grouped_images[group_identifier].append(image_file)

        for group_identifier, group_images in grouped_images.items():
            slide = ppt.slides.add_slide(layout)
            num_images = len(group_images)

            # 슬라이드 제목 설정 (두 번째 상위 폴더 이름으로)
            folder_name = os.path.basename(os.path.dirname(group_images[0]))  # 상위 폴더 이름 추출
            grandparent_folder_name = os.path.basename(os.path.dirname(os.path.dirname(group_images[0])))  # 두 번째 상위 폴더 이름 추출
            title_shape = slide.shapes.title
            title_shape.text = grandparent_folder_name  # 슬라이드 제목으로 두 번째 상위 폴더 이름 설정

            # 슬라이드 배경색을 검은색으로 설정
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # 검은색 (RGB 값: 0, 0, 0)

            # 폰트 설정
            title_text_frame = title_shape.text_frame
            title_paragraph = title_text_frame.paragraphs[0]
            title_run = title_paragraph.runs[0]
            title_font = title_run.font
            title_font.size = Pt(14)  # 폰트 크기를 작게 설정

            # 텍스트 색상을 흰색으로 변경
            title_font.color.rgb = RGBColor(255, 255, 255)  # 흰색 (RGB 값: 255, 255, 255)

            # 텍스트 상자의 위치 조정 (오른쪽 상단)
            title_left = Inches(6.5)  # 좌측 여백 (오른쪽으로 이동)
            title_top = Inches(0.5)   # 상단 여백
            title_width = Inches(3)   # 텍스트 상자 가로 크기
            title_height = Inches(0.0)  # 텍스트 상자 높이
            title_shape.left = title_left
            title_shape.top = title_top
            title_shape.width = title_width
            title_shape.height = title_height
            title_paragraph.alignment = PP_ALIGN.RIGHT  # 텍스트를 오른쪽 정렬

            # "io"로 시작하는 그룹에 대한 처리
            if group_identifier.startswith("io"):
                slide_width = Inches(8.34)  # 이미지 너비 (슬라이드 너비)
                slide_height = Inches(5.56)  # 이미지 높이 (슬라이드 높이)
                img_width = slide_width / 2.8  # 이미지 너비 (4분할)
                img_height = img_width * 0.67  # 이미지 높이를 이미지 너비의 1.5배로 조절

                # 파일 이름에서 위치 정보 추출 (예: "io-01"에서 "01"을 추출)
                positions = [int(os.path.splitext(os.path.basename(image_file))[0].split("-")[1]) for image_file in group_images]

                # 위치 정보를 기록한 딕셔너리를 사용하여 이미지 위치 설정
                position_mapping = {
                    "io-01": (1, 0),  # io-01을 4 위치에 배치
                    "io-02": (1, 1),  # io-02를 5 위치에 배치
                    "io-03": (1, 2),  # io-03을 6 위치에 배치
                    "io-04": (0, 1),  # io-04를 2 위치에 배치
                    "io-05": (2, 1),  # io-05를 8 위치에 배치
                    "io-06": (2, 0),  # io-06을 1 위치에 배치
                    "io-07": (2, 2),  # io-07을 3 위치에 배치
                }

                # 여백 설정
                top_margin = Inches(0.7)  # 상단 여백
                left_margin = Inches(0.5)  # 좌측 여백

                for image_file, (row, col) in zip(group_images, position_mapping.values()):
                    left = col * img_width + left_margin
                    top = row * img_height + top_margin
                    pic = slide.shapes.add_picture(image_file, left, top, img_width, img_height)

                    # 이미지 크기 조절 (이미지 높이와 너비에 맞게)
                    adjust_image_size(pic, img_width, img_height)

                    pic.width = int(Inches(3.0))
                    pic.height = int(Inches(2.0))

       
            # "eo"로 시작하는 그룹에 대한 처리
            elif group_identifier.startswith("eo"):
                # "eo"로 시작하는 그룹을 2장씩 나누어 처리
                num_slides = (num_images + 1) // 2  # 이미지를 2장씩 나누므로 슬라이드 개수 계산 (나머지 이미지를 위해 1을 더함)

                if num_slides > 0:  # num_slides가 0보다 큰 경우에만 슬라이드 생성
                    for slide_index in range(num_slides):
                        slide = ppt.slides.add_slide(layout)  # 새로운 슬라이드 생성

                        # 슬라이드의 제목 설정 (두 번째 상위 폴더 이름으로, eo 포함)
                        folder_name = os.path.basename(os.path.dirname(group_images[0]))  # 상위 폴더 이름 추출
                        grandparent_folder_name = os.path.basename(os.path.dirname(os.path.dirname(group_images[0])))  # 두 번째 상위 폴더 이름 추출
                        title_shape = slide.shapes.title
                        title_text = grandparent_folder_name  # 두 번째 상위 폴더 이름과 상위 폴더 이름을 결합
                        title_shape.text = title_text  # 슬라이드 제목으로 설정

                        # 폰트 설정
                        title_text_frame = title_shape.text_frame
                        title_paragraph = title_text_frame.paragraphs[0]
                        title_run = title_paragraph.runs[0]
                        title_font = title_run.font
                        title_font.size = Pt(14)  # 폰트 크기를 작게 설정

                        # 텍스트 색상을 흰색으로 변경
                        title_font.color.rgb = RGBColor(255, 255, 255)  # 흰색 (RGB 값: 255, 255, 255)

                        # 슬라이드 배경색을 검은색으로 설정
                        background = slide.background
                        fill = background.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(0, 0, 0)  # 검은색 (RGB 값: 0, 0, 0)

                        # 텍스트 상자의 위치 조정 (오른쪽 상단)
                        title_left = Inches(6.5)  # 좌측 여백 (오른쪽으로 이동)
                        title_top = Inches(0.5)   # 상단 여백
                        title_width = Inches(3)   # 텍스트 상자 가로 크기
                        title_height = Inches(0.0)  # 텍스트 상자 높이
                        title_shape.left = title_left
                        title_shape.top = title_top
                        title_shape.width = title_width
                        title_shape.height = title_height
                        title_paragraph.alignment = PP_ALIGN.RIGHT  # 텍스트를 오른쪽 정렬

                        slide_width = int(Inches(10))  # 이미지 너비 (슬라이드 너비)
                        slide_height = int(Inches(7.5))  # 이미지 높이 (슬라이드 높이)
                        img_width = slide_width / 2  # 이미지 너비 (2분할)
                        
                        # 현재 슬라이드에 2장의 이미지 추가
                        for i in range(2):
                            image_index = slide_index * 2 + i
                            if image_index < num_images:
                                left = i * img_width + Inches(0)  # 좌측 여백 추가
                                top = Inches(2.1)  # 상단 여백 추가
                                pic = slide.shapes.add_picture(group_images[image_index], left, top, img_width, slide_height)

                                

                                # 이미지 크기 조절 (슬라이드에 맞게)
                                adjust_image_size(pic, img_width, slide_height)

                                pic.width = int(Inches(5.0))
                                pic.height = int(Inches(3.3))

            elif group_identifier.startswith("lateral"):
                for _ in range(2):
                    ppt.slides.add_slide(layout)

                slide_width = ppt.slide_width
                slide_height = ppt.slide_height
                img_width = slide_width / num_images

                for i, image_file in enumerate(group_images):
                    left = i * img_width
                    top = 0
                    pic = slide.shapes.add_picture(image_file, left, top, img_width, slide_height)

                    # 이미지 크기 조절 (슬라이드에 맞게)
                    adjust_image_size(pic, img_width, slide_height)

                    # 이미지 크기를 0.8로 줄임
                    pic.width = int(Inches(5.65))
                    pic.height = int(Inches(6.8))

                    # 이미지 가운데 정렬
                    pic.left = int((slide_width - pic.width) / 2)
                    pic.top = int((slide_height - pic.height) / 2)

            elif group_identifier.endswith("pa"):
                for _ in range(1):
                    ppt.slides.add_slide(layout)

                slide_width = ppt.slide_width
                slide_height = ppt.slide_height
                img_width = slide_width / num_images

                for i, image_file in enumerate(group_images):
                    left = i * img_width
                    top = 0
                    pic = slide.shapes.add_picture(image_file, left, top, img_width, slide_height)

                    # 이미지 크기 조절 (슬라이드에 맞게)
                    adjust_image_size(pic, img_width, slide_height)

                    # 이미지 크기를 0.8로 줄임
                    pic.width = int(Inches(6.6))
                    pic.height = int(Inches(6.6))

                    # 이미지 가운데 정렬
                    pic.left = int((slide_width - pic.width) / 2)
                    pic.top = int((slide_height - pic.height) / 2)

            elif group_identifier.endswith("no"):

                slide_width = ppt.slide_width
                slide_height = ppt.slide_height
                img_width = slide_width / num_images

                for i, image_file in enumerate(group_images):
                    left = i * img_width
                    top = 0
                    pic = slide.shapes.add_picture(image_file, left, top, img_width, slide_height)

                    # 이미지 크기 조절 (슬라이드에 맞게)
                    adjust_image_size(pic, img_width, slide_height)

                    # 이미지 크기를 0.8로 줄임
                    pic.width = int(Inches(10.0))
                    pic.height = int(Inches(4.93))

                    # 이미지 가운데 정렬
                    pic.left = int((slide_width - pic.width) / 2)
                    pic.top = int((slide_height - pic.height) / 2)

            elif group_identifier.endswith("wrist"):
                
                slide_width = ppt.slide_width
                slide_height = ppt.slide_height
                img_width = slide_width / num_images

                for i, image_file in enumerate(group_images):
                    left = i * img_width
                    top = 0
                    pic = slide.shapes.add_picture(image_file, left, top, img_width, slide_height)

                    # 이미지 크기 조절 (슬라이드에 맞게)
                    adjust_image_size(pic, img_width, slide_height)

                    # 이미지 크기를 0.8로 줄임
                    pic.width = int(Inches(5.92))
                    pic.height = int(Inches(7.2))

                    # 이미지 가운데 정렬
                    pic.left = int((slide_width - pic.width) / 2)
                    pic.top = int((slide_height - pic.height) / 2)

        



            # "eo", "io"로 시작하지 않는 그룹에 대한 처리
            else:
                slide_width = ppt.slide_width
                slide_height = ppt.slide_height
                img_width = slide_width / num_images

                for i, image_file in enumerate(group_images):
                    left = i * img_width
                    top = 0
                    pic = slide.shapes.add_picture(image_file, left, top, img_width, slide_height)

                    # 이미지 크기 조절 (슬라이드에 맞게)
                    adjust_image_size(pic, img_width, slide_height)

                    # 이미지 크기를 0.8로 줄임
                    pic.width = int(pic.width * 0.8)
                    pic.height = int(pic.height * 0.8)

                    # 이미지 가운데 정렬
                    pic.left = int((slide_width - pic.width) / 2)
                    pic.top = int((slide_height - pic.height) / 2)


    # 1번째 슬라이드에 이미지 추가
    if len(ppt.slides) >= 1:
        first_slide = ppt.slides[0]  # 1번째 슬라이드 선택

        # 이미지 슬라이드 레이아웃(5)으로 변경
        layout = ppt.slide_layouts[5]
        first_slide.layout = layout

        background = first_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(253, 249, 236)

        # 이미지 파일 경로
        image1_path = resource_path("ii.png")  # 첫 번째 이미지 경로
        image2_path = resource_path("jj.png")  # 두 번째 이미지 경로

        # 첫 번째 이미지 크기 및 위치 설정
        slide_width = Inches(10.0)  # 이미지 너비 (슬라이드 너비)
        slide_height = Inches(7.5)  # 이미지 높이 (슬라이드 높이)
        img_width1 = slide_width   # 첫 번째 이미지 너비
        img_height1 = slide_height  # 첫 번째 이미지 높이
        left1 = int(Inches(0))  # 첫 번째 이미지를 원하는 위치로 이동 (Inches를 조절하여 위치 조정)
        top1 = int(Inches(0))
        pic1 = first_slide.shapes.add_picture(image1_path, left1, top1, img_width1, img_height1)

        # 두 번째 이미지 크기 및 위치 설정
        img_width2 = slide_width/2   # 두 번째 이미지 너비
        img_height2 = slide_height/5  # 두 번째 이미지 높이
        left2 = int(Inches(5))  # 두 번째 이미지를 원하는 위치로 이동 (Inches를 조절하여 위치 조정)
        top2 = int(Inches(5.8))
        pic2 = first_slide.shapes.add_picture(image2_path, left2, top2, img_width2, img_height2)


        # 이미지 크기 조절 (슬라이드에 맞게)
        adjust_image_size(pic1, img_width1, img_height1)
        adjust_image_size(pic2, img_width2, img_height2)

        # 추가 텍스트 상자 추가 (인치 단위로 위치 조정)
        left = Inches(0.8)    # 텍스트 상자 왼쪽 여백
        top = Inches(1.4)     # 텍스트 상자 상단 여백 (인치로 조정)
        width = Inches(8)   # 텍스트 상자 너비
        height = Inches(2)  # 텍스트 상자 높이

        text_box = first_slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame


        # 첫 번째 문단 (Skeletal class II)
        p1 = text_frame.add_paragraph()
        p1.text = "정밀 진단 분석 및 치료계획"
        p1.font.size = Pt(49)  # 폰트 크기 적용
        p1.font.bold = True  # 폰트 굵게 설정
     
        p2 = text_frame.add_paragraph()
        p2.text = ""
        p2.font.size = Pt(30)  # 폰트 크기 적용
 
        p3 = text_frame.add_paragraph()
        p3.text = ""
        p3.font.size = Pt(30)  # 폰트 크기 적용
    
        p4 = text_frame.add_paragraph()
        p4.text = "OOO"
        p4.font.size = Pt(30)  # 폰트 크기 적용
        p4.font.color.rgb = RGBColor(136, 136, 136)  
        p4.alignment = PP_ALIGN.CENTER
      
        p5 = text_frame.add_paragraph()
        p5.text = "DOB : 0000. 00. 00"
        p5.font.size = Pt(30)  # 폰트 크기 적용
        p5.font.color.rgb = RGBColor(136, 136, 136)  
        p5.alignment = PP_ALIGN.CENTER

        p6 = text_frame.add_paragraph()
        p6.text = "00Y 0M / F"
        p6.font.size = Pt(30)  # 폰트 크기 적용
        p6.font.color.rgb = RGBColor(136, 136, 136)  
        p6.alignment = PP_ALIGN.CENTER
        
        p7 = text_frame.add_paragraph()
        p7.text = "Dx. Day : 0000. 00. 00 "
        p7.font.size = Pt(30)  # 폰트 크기 적용
        p7.font.color.rgb = RGBColor(136, 136, 136)
        p7.alignment = PP_ALIGN.CENTER

    # 6번째 슬라이드에 이미지 추가
    if len(ppt.slides) >= 6:
        sixth_slide = ppt.slides[5]  # 6번째 슬라이드 선택

        # 이미지 슬라이드 레이아웃(5)으로 변경
        layout = ppt.slide_layouts[5]
        sixth_slide.layout = layout

        # 이미지 파일 경로
        image1_path = resource_path("bb.png")  # 첫 번째 이미지 경로
        image2_path = resource_path("cc.png")  # 두 번째 이미지 경로

        # 첫 번째 이미지 크기 및 위치 설정
       
        img_width1 = Inches(5.0)  # 첫 번째 이미지 너비
        img_height1 = Inches(4.7)  # 첫 번째 이미지 높이
        left1 = int(Inches(0.0))  # 첫 번째 이미지를 원하는 위치로 이동 (Inches를 조절하여 위치 조정)
        top1 = int(Inches(1.3))
        pic1 = sixth_slide.shapes.add_picture(image1_path, left1, top1, img_width1, img_height1)

        # 두 번째 이미지 크기 및 위치 설정
        img_width2 = Inches(5.0)   # 두 번째 이미지 너비
        img_height2 = Inches(6.8)  # 두 번째 이미지 높이
        left2 = int(Inches(5.0))  # 두 번째 이미지를 원하는 위치로 이동 (Inches를 조절하여 위치 조정)
        top2 = int(Inches(0.7))
        pic2 = sixth_slide.shapes.add_picture(image2_path, left2, top2, img_width2, img_height2)

        # 슬라이드 배경색을 검은색으로 설정
        for slide in ppt.slides:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # 검은색 배경

        # 이미지 크기 조절 (슬라이드에 맞게)
        adjust_image_size(pic1, img_width1, img_height1)
        adjust_image_size(pic2, img_width2, img_height2)

        # 제목 텍스트 추가
        title = sixth_slide.shapes.title
        title.text = " "
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

    

        

        




    # 7번째 슬라이드에 이미지 추가
    if len(ppt.slides) >= 7:
        seventh_slide = ppt.slides[6]  # 7번째 슬라이드 선택

        layout = ppt.slide_layouts[5]
        seventh_slide.layout = layout

        # 이미지 파일 경로
        image_path = resource_path("dd.png")  # 이미지 경로

        # 이미지 크기 및 위치 설정
       
        img_width = Inches(6.8)    # 이미지 너비
        img_height = Inches(5.8)   # 이미지 높이
        left = int(Inches(1.6))  # 첫 번째 이미지를 원하는 위치로 이동 (Inches를 조절하여 위치 조정)
        top = int(Inches(0.85))
        pic = seventh_slide.shapes.add_picture(image_path, left, top, img_width, img_height)

        # 슬라이드 배경색을 검은색으로 설정
        for slide in ppt.slides:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # 검은색 배경

        # 이미지 크기 조절 (슬라이드에 맞게)
        adjust_image_size(pic, img_width, img_height)

        # 제목 텍스트 추가
        title = seventh_slide.shapes.title
        title.text = " "
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨


    # 9번째 슬라이드에 이미지 추가
    if len(ppt.slides) >= 9:
        seventh_slide = ppt.slides[8]  # 9번째 슬라이드 선택

        layout = ppt.slide_layouts[5]
        seventh_slide.layout = layout
        for _ in range(6):
                    ppt.slides.add_slide(layout)
        
        # 이미지 파일 경로
        image_path = resource_path("ee.png")  # 이미지 경로

        img_width = Inches(7.4)    # 이미지 너비
        img_height = Inches(6.6)   # 이미지 높이
        left = int(Inches(1.4))  # 첫 번째 이미지를 원하는 위치로 이동 (Inches를 조절하여 위치 조정)
        top = int(Inches(0.5))
        pic = seventh_slide.shapes.add_picture(image_path, left, top, img_width, img_height)

        # 슬라이드 배경색을 검은색으로 설정
        for slide in ppt.slides:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # 검은색 배경


        # 이미지 크기 조절 (슬라이드에 맞게)
        adjust_image_size(pic, img_width, img_height)

        # 제목 텍스트 추가
        title = seventh_slide.shapes.title
        title.text = " "
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨


    # 12번째 슬라이드에 이미지 추가
    if len(ppt.slides) >= 12:
        seventh_slide = ppt.slides[11]  # 12번째 슬라이드 선택

        layout = ppt.slide_layouts[5]
        seventh_slide.layout = layout


        # 이미지 파일 경로
        image1_path = resource_path("ee1.png")  # 이미지 경로
        image2_path = resource_path("ee2.png")  # 이미지 경로
        image3_path = resource_path("ee3.png")  # 이미지 경로
        image4_path = resource_path("ee4.png")  # 이미지 경로
        image5_path = resource_path("ee5.png")  # 이미지 경로
        image6_path = resource_path("ee6.png")  # 이미지 경로
        image7_path = resource_path("ee7.png")  # 이미지 경로
        image8_path = resource_path("ee8.jpg")  # 이미지 경로


        # 첫 번째 이미지 크기 및 위치 설정
        
        img_width1 = Inches(4.0)   # 첫 번째 이미지 너비
        img_height1 = Inches(2.52)  # 첫 번째 이미지 높이
        left1 = int(Inches(0.2))  # 첫 번째 이미지를 원하는 위치로 이동 (Inches를 조절하여 위치 조정)
        top1 = int(Inches(1.3))
        pic1 = seventh_slide.shapes.add_picture(image1_path, left1, top1, img_width1, img_height1)

        # 두 번째 이미지 크기 및 위치 설정
        img_width2 = Inches(6.0)   # 두 번째 이미지 너비
        img_height2 = Inches(3.2)  # 두 번째 이미지 높이
        left2 = int(Inches(0.2))  # 두 번째 이미지를 원하는 위치로 이동 (Inches를 조절하여 위치 조정)
        top2 = int(Inches(4.0))
        pic2 = seventh_slide.shapes.add_picture(image2_path, left2, top2, img_width2, img_height2)

        # 세 번째 이미지 크기 및 위치 설정
        img_width3 = Inches(2.6)    # 세 번째 이미지 너비
        img_height3 = Inches(1.3)  # 세 번째 이미지 높이
        left3 = int(Inches(7.0))  # 세 번째 이미지를 원하는 위치로 이동 (Inches를 조절하여 위치 조정)
        top3 = int(Inches(0.2))
        pic3 = seventh_slide.shapes.add_picture(image3_path, left3, top3, img_width3, img_height3)

        # 네 번째 이미지 크기 및 위치 설정
        img_width4 = Inches(3.3)    # 네 번째 이미지 너비
        img_height4 = Inches(1.3)  # 네 번째 이미지 높이
        left4 = int(Inches(6.5))  # 네 번째 이미지를 원하는 위치로 이동 (Inches를 조절하여 위치 조정)
        top4 = int(Inches(1.55))
        pic4 = seventh_slide.shapes.add_picture(image4_path, left4, top4, img_width4, img_height4)

        # 다섯 번째 이미지 크기 및 위치 설정
        img_width5 = Inches(3.3)  # 다섯 번째 이미지 너비
        img_height5 = Inches(1.3) # 다섯 번째 이미지 높이
        left5 = int(Inches(6.5))  # 다섯 번째 이미지를 원하는 위치로 이동 (Inches를 조절하여 위치 조정)
        top5 = int(Inches(2.9))
        pic5 = seventh_slide.shapes.add_picture(image5_path, left5, top5, img_width5, img_height5)

        # 여섯 번째 이미지 크기 및 위치 설정
        img_width6 = Inches(3.3)  
        img_height6 = Inches(1.3)
        left6 = int(Inches(6.5))  
        top6 = int(Inches(4.25))
        pic6 = seventh_slide.shapes.add_picture(image6_path, left6, top6, img_width6, img_height6)

        # 일곱 번째 이미지 크기 및 위치 설정
        img_width7 = Inches(3.3)  
        img_height7 = Inches(1.3) 
        left7 = int(Inches(6.5))  
        top7 = int(Inches(5.6))
        pic7 = seventh_slide.shapes.add_picture(image7_path, left7, top7, img_width7, img_height7)

        # 여덟 번째 이미지 크기 및 위치 설정
        img_width8 = Inches(3.0)  
        img_height8 = Inches(3.8) 
        left8 = int(Inches(3.7))  
        top8 = int(Inches(3.6))
        pic8 = seventh_slide.shapes.add_picture(image8_path, left8, top8, img_width8, img_height8)

        # 슬라이드 배경색을 검은색으로 설정
        for slide in ppt.slides:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # 검은색 배경

        # 이미지 크기 조절 (슬라이드에 맞게)
        adjust_image_size(pic1, img_width1, img_height1)
        adjust_image_size(pic2, img_width2, img_height2)
        adjust_image_size(pic3, img_width3, img_height3)
        adjust_image_size(pic4, img_width4, img_height4)
        adjust_image_size(pic5, img_width5, img_height5)
        adjust_image_size(pic6, img_width6, img_height6)
        adjust_image_size(pic7, img_width7, img_height7)
        adjust_image_size(pic8, img_width8, img_height8)


        # 제목 텍스트 추가
        title = seventh_slide.shapes.title
        title.text = "HW"
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨



    # 13번째 슬라이드 선택
    if len(ppt.slides) >= 13:
        thirteenth_slide = ppt.slides[12]  # 13번째 슬라이드 선택 (0부터 시작하므로 12번째)

        # 제목 텍스트 추가
        title = thirteenth_slide.shapes.title
        title.text = "Diagnosis"
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 추가 텍스트 상자 추가 (인치 단위로 위치 조정)
        left = Inches(0.8)    # 텍스트 상자 왼쪽 여백
        top = Inches(1.4)     # 텍스트 상자 상단 여백 (인치로 조정)
        width = Inches(8)   # 텍스트 상자 너비
        height = Inches(2)  # 텍스트 상자 높이

        text_box = thirteenth_slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame

        # 텍스트 크기를 키워줄 폰트 크기
        font_size = Pt(20)

        # 첫 번째 문단 (Skeletal class II)
        p1 = text_frame.add_paragraph()
        p1.text = "memo"
        p1.font.size = font_size  # 폰트 크기 적용
        p1.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # p1 = text_frame.add_paragraph()
        # p1.text = ""
        # p1.font.size = font_size  # 폰트 크기 적용
        # p1.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # # 두 번째 문단 (Dental class II (both))
        # p2 = text_frame.add_paragraph()
        # p2.text = "Dental class II (both)"
        # p2.font.size = font_size  # 폰트 크기 적용
        # p2.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # # 세 번째 문단 (congenital missing : #23)
        # p3 = text_frame.add_paragraph()
        # p3.text = "congenital missing : #23"
        # p3.font.size = font_size  # 폰트 크기 적용
        # p3.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # # 네 번째 문단 (- prolonged retention: #63)
        # p4 = text_frame.add_paragraph()
        # p4.text = "- prolonged retention: #63"
        # p4.font.size = font_size  # 폰트 크기 적용
        # p4.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # # 다섯 번째 문단 (upper midline deviation to Lt.)
        # p5 = text_frame.add_paragraph()
        # p5.text = "upper midline deviation to Lt."
        # p5.font.size = font_size  # 폰트 크기 적용
        # p5.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # # 여섯 번째 문단 (lower crowding)
        # p6 = text_frame.add_paragraph()
        # p6.text = "lower crowding"
        # p6.font.size = font_size  # 폰트 크기 적용
        # p6.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

    # 14번째 슬라이드 선택
    if len(ppt.slides) >= 14:
        thirteenth_slide = ppt.slides[13]  # 14번째 슬라이드 선택 (0부터 시작하므로 12번째)

        # 이미지 파일 경로
        image_path = resource_path("gg.png")  # 이미지 경로

        # 이미지 크기 및 위치 설정
        slide_width = Inches(2.5)  # 이미지 너비 (슬라이드 너비)
        slide_height = Inches(1.5)  # 이미지 높이 (슬라이드 높이)
        img_width = slide_width   # 이미지 너비
        img_height = slide_height  # 이미지 높이
        left = int(Inches(7.0))  # 첫 번째 이미지를 원하는 위치로 이동 (Inches를 조절하여 위치 조정)
        top = int(Inches(0.2))
        pic = thirteenth_slide.shapes.add_picture(image_path, left, top, img_width, img_height)



        # 이미지 크기 조절 (슬라이드에 맞게)
        adjust_image_size(pic, img_width, img_height)

        # 제목 텍스트 추가
        title = thirteenth_slide.shapes.title
        title.text = "Tx. plan"
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 추가 텍스트 상자 추가 (인치 단위로 위치 조정)
        left = Inches(0.8)    # 텍스트 상자 왼쪽 여백
        top = Inches(1.4)     # 텍스트 상자 상단 여백 (인치로 조정)
        width = Inches(8)   # 텍스트 상자 너비
        height = Inches(2)  # 텍스트 상자 높이

        text_box = thirteenth_slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame

       

        # 첫 번째 문단 (Skeletal class II)
        p1 = text_frame.add_paragraph()
        p1.text = "#make a note"
        p1.font.size = Pt(20)  # 폰트 크기 적용
        p1.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨
        p1.font.bold = True  # 폰트 굵게 설정

        p1 = text_frame.add_paragraph()
        p1.text = ""
        p1.font.size = Pt(20)  # 폰트 크기 적용
        p1.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨
        p1.font.bold = True  # 폰트 굵게 설정

        p1 = text_frame.add_paragraph()
        p1.text = "Make a note"
        p1.font.size = Pt(16)  # 폰트 크기 적용
        p1.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 두 번째 문단 (Dental class II (both))
        p2 = text_frame.add_paragraph()
        p2.text = ""
        p2.font.size = Pt(16)  # 폰트 크기 적용
        p2.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 세 번째 문단 (congenital missing : #23)
        p3 = text_frame.add_paragraph()
        p3.text = ""
        p3.font.size = Pt(16)  # 폰트 크기 적용
        p3.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 네 번째 문단 (- prolonged retention: #63)
        p4 = text_frame.add_paragraph()
        p4.text = ""
        p4.font.size = Pt(16)  # 폰트 크기 적용
        p4.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 다섯 번째 문단 (upper midline deviation to Lt.)
        p5 = text_frame.add_paragraph()
        p5.text = ""
        p5.font.size = Pt(16)  # 폰트 크기 적용
        p5.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 여섯 번째 문단 (lower crowding)
        p6 = text_frame.add_paragraph()
        p6.text = ""
        p6.font.size = Pt(16)  # 폰트 크기 적용
        p6.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 여섯 번째 문단 (lower crowding)
        p8 = text_frame.add_paragraph()
        p8.text = "* 상악 정중선 좌측으로 약간의 변이 가능성"
        p8.font.size = Pt(14)  # 폰트 크기 적용
        p8.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 여섯 번째 문단 (lower crowding)
        p9 = text_frame.add_paragraph()
        p9.text = "* Miniscrew 식립 필요성, IPR 가능성"
        p9.font.size = Pt(14)  # 폰트 크기 적용
        p9.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 여섯 번째 문단 (lower crowding)
        p10 = text_frame.add_paragraph()
        p10.text = "* 치근흡수, 치은퇴축, Black Triangle 발생 가능성"
        p10.font.size = Pt(14)  # 폰트 크기 적용
        p10.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 여섯 번째 문단 (lower crowding)
        p11 = text_frame.add_paragraph()
        p11.text = "* #24 Palatal cusp OE 필요성"
        p11.font.size = Pt(14)  # 폰트 크기 적용
        p11.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        p12 = text_frame.add_paragraph()
        p12.text = "* 장치 착용 시간 매우 중요함 (하루 20시간 이상, 식사/양치 시간 제외)"
        p12.font.size = Pt(14)  # 폰트 크기 적용
        p12.font.color.rgb = RGBColor(255, 0, 0)  # 흰색 글씨
        p12.font.bold = True  # 폰트 굵게 설정

        p13 = text_frame.add_paragraph()
        p13.text = "* 고무줄 협조도 중요함. "
        p13.font.size = Pt(14)  # 폰트 크기 적용
        p13.font.color.rgb = RGBColor(255, 0, 0)  # 흰색 글씨
        p13.font.bold = True  # 폰트 굵게 설정

        


        # 이미지 크기 조절 (슬라이드에 맞게)
        adjust_image_size(pic, img_width, img_height)

        


    # 15번째 슬라이드 선택
    if len(ppt.slides) >= 15:
        thirteenth_slide = ppt.slides[14]  # 14번째 슬라이드 선택 (0부터 시작하므로 12번째)

        # 제목 텍스트 추가
        title = thirteenth_slide.shapes.title
        title.text = "Tx. plan"
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 추가 텍스트 상자 추가 (인치 단위로 위치 조정)
        left = Inches(0.8)    # 텍스트 상자 왼쪽 여백
        top = Inches(1.4)     # 텍스트 상자 상단 여백 (인치로 조정)
        width = Inches(8)   # 텍스트 상자 너비
        height = Inches(2)  # 텍스트 상자 높이

        text_box = thirteenth_slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame

       

        # 첫 번째 문단 (Skeletal class II)
        p1 = text_frame.add_paragraph()
        p1.text = "#make a note"
        p1.font.size = Pt(20)  # 폰트 크기 적용
        p1.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨
        p1.font.bold = True  # 폰트 굵게 설정

        p1 = text_frame.add_paragraph()
        p1.text = ""
        p1.font.size = Pt(20)  # 폰트 크기 적용
        p1.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨
        p1.font.bold = True  # 폰트 굵게 설정

        p1 = text_frame.add_paragraph()
        p1.text = "Make a note"
        p1.font.size = Pt(16)  # 폰트 크기 적용
        p1.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 두 번째 문단 (Dental class II (both))
        p2 = text_frame.add_paragraph()
        p2.text = ""
        p2.font.size = Pt(16)  # 폰트 크기 적용
        p2.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 세 번째 문단 (congenital missing : #23)
        p3 = text_frame.add_paragraph()
        p3.text = ""
        p3.font.size = Pt(16)  # 폰트 크기 적용
        p3.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 네 번째 문단 (- prolonged retention: #63)
        p4 = text_frame.add_paragraph()
        p4.text = ""
        p4.font.size = Pt(16)  # 폰트 크기 적용
        p4.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 다섯 번째 문단 (upper midline deviation to Lt.)
        p5 = text_frame.add_paragraph()
        p5.text = ""
        p5.font.size = Pt(16)  # 폰트 크기 적용
        p5.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 여섯 번째 문단 (lower crowding)
        p6 = text_frame.add_paragraph()
        p6.text = ""
        p6.font.size = Pt(16)  # 폰트 크기 적용
        p6.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 여섯 번째 문단 (lower crowding)
        p7 = text_frame.add_paragraph()
        p7.text = ""
        p7.font.size = Pt(16)  # 폰트 크기 적용
        p7.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 여섯 번째 문단 (lower crowding)
        p8 = text_frame.add_paragraph()
        p8.text = ""
        p8.font.size = Pt(16)  # 폰트 크기 적용
        p8.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 여섯 번째 문단 (lower crowding)
        p9 = text_frame.add_paragraph()
        p9.text = ""
        p9.font.size = Pt(16)  # 폰트 크기 적용
        p9.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 여섯 번째 문단 (lower crowding)
        p10 = text_frame.add_paragraph()
        p10.text = ""
        p10.font.size = Pt(16)  # 폰트 크기 적용
        p10.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 여섯 번째 문단 (lower crowding)
        p11 = text_frame.add_paragraph()
        p11.text = ""
        p11.font.size = Pt(16)  # 폰트 크기 적용
        p11.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        p12 = text_frame.add_paragraph()
        p12.text = "* 상악 정중선 좌측으로 약간의 변이 가능성"
        p12.font.size = Pt(14)  # 폰트 크기 적용
        p12.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        p13 = text_frame.add_paragraph()
        p13.text = "* Miniscrew 식립 필요성, IPR 가능성"
        p13.font.size = Pt(14)  # 폰트 크기 적용
        p13.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        p14 = text_frame.add_paragraph()
        p14.text = "* 치근흡수, 치은퇴축, Black Triangle 발생 가능성"
        p14.font.size = Pt(14)  # 폰트 크기 적용
        p14.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        p15 = text_frame.add_paragraph()
        p15.text = "* #24 Palatal cusp OE 필요성"
        p15.font.size = Pt(14)  # 폰트 크기 적용
        p15.font.color.rgb = RGBColor(255, 255, 255)   # 빨간색 글씨
        
        p17 = text_frame.add_paragraph()
        p17.text = "* 고무줄 협조도 중요함. "
        p17.font.size = Pt(14)  # 폰트 크기 적용
        p17.font.color.rgb = RGBColor(255, 0, 0)   # 빨간색 글씨
        p17.font.bold = True  # 폰트 굵게 설정


        # 16번째 슬라이드 선택
    if len(ppt.slides) >= 16:
        sixteenth_slide = ppt.slides[15]  # 14번째 슬라이드 선택 (0부터 시작하므로 12번째)

        # 이미지 파일 경로
        image_path = resource_path("hh.png")  # 이미지 경로

        # 이미지 크기 및 위치 설정
        slide_width = Inches(4)  # 이미지 너비 (슬라이드 너비)
        slide_height = Inches(4)  # 이미지 높이 (슬라이드 높이)
        img_width = slide_width   # 이미지 너비
        img_height = slide_height  # 이미지 높이
        left = int(Inches(5.5))  # 첫 번째 이미지를 원하는 위치로 이동 (Inches를 조절하여 위치 조정)
        top = int(Inches(3.0))
        pic = sixteenth_slide.shapes.add_picture(image_path, left, top, img_width, img_height)



        # 이미지 크기 조절 (슬라이드에 맞게)
        adjust_image_size(pic, img_width, img_height)

        # 제목 텍스트 추가
        title = sixteenth_slide.shapes.title
        title.text = "교정 치료 시 주의사항"
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 추가 텍스트 상자 추가 (인치 단위로 위치 조정)
        left = Inches(0.8)    # 텍스트 상자 왼쪽 여백
        top = Inches(1.4)     # 텍스트 상자 상단 여백 (인치로 조정)
        width = Inches(8)   # 텍스트 상자 너비
        height = Inches(2)  # 텍스트 상자 높이

        text_box = sixteenth_slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame

       # 첫 번째 문단 (Dental class II (both))
        p1 = text_frame.add_paragraph()
        p1.text = "1. 양치는 깨끗이 ! ★ ★ ★"
        p1.font.size = Pt(20)  # 폰트 크기 적용
        p1.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 두 번째 문단 (Dental class II (both))
        p2 = text_frame.add_paragraph()
        p2.text = "   - 치간칫솔 → 칫솔질 및 치실 "
        p2.font.size = Pt(20)  # 폰트 크기 적용
        p2.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 세 번째 문단 (congenital missing : #23)
        p3 = text_frame.add_paragraph()
        p3.text = "   - 장치 주위로 하얀 충치 발생 가능"
        p3.font.size = Pt(20)  # 폰트 크기 적용
        p3.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 네 번째 문단 (- prolonged retention: #63)
        p4 = text_frame.add_paragraph()
        p4.text = "2. 발생 가능한 부작용"
        p4.font.size = Pt(20)  # 폰트 크기 적용
        p4.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 다섯 번째 문단 (upper midline deviation to Lt.)
        p5 = text_frame.add_paragraph()
        p5.text = "   - 치아 탈회 (충치의 발생) "
        p5.font.size = Pt(20)  # 폰트 크기 적용
        p5.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 여섯 번째 문단 (lower crowding)
        p6 = text_frame.add_paragraph()
        p6.text = "   - 치근 흡수의 가능성 "
        p6.font.size = Pt(20)  # 폰트 크기 적용
        p6.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 여섯 번째 문단 (lower crowding)
        p7 = text_frame.add_paragraph()
        p7.text = "   - Black triangle (▲) "
        p7.font.size = Pt(20)  # 폰트 크기 적용
        p7.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨

        # 여섯 번째 문단 (lower crowding)
        p8 = text_frame.add_paragraph()
        p8.text = "   - 일시적 치아 변색 "
        p8.font.size = Pt(20)  # 폰트 크기 적용
        p8.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 글씨
        


        # 이미지 크기 조절 (슬라이드에 맞게)
        adjust_image_size(pic, img_width, img_height)




    # 17번째 슬라이드에 이미지 추가
    if len(ppt.slides) >= 17:
        seventh_slide = ppt.slides[16]  # 9번째 슬라이드 선택

        layout = ppt.slide_layouts[6]
        seventh_slide.layout = layout
        # 이미지 파일 경로
        image_path = resource_path("ff.png")  # 이미지 경로

        # 이미지 크기 및 위치 설정
        slide_width = Inches(10)  # 이미지 너비 (슬라이드 너비)
        slide_height = Inches(7.5)  # 이미지 높이 (슬라이드 높이)
        img_width = slide_width   # 이미지 너비
        img_height = slide_height  # 이미지 높이
        left = int(Inches(0))  # 첫 번째 이미지를 원하는 위치로 이동 (Inches를 조절하여 위치 조정)
        top = int(Inches(0))
        pic = seventh_slide.shapes.add_picture(image_path, left, top, img_width, img_height)

        # 슬라이드 배경색을 검은색으로 설정
        for slide in ppt.slides:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # 검은색 배경


        # 이미지 크기 조절 (슬라이드에 맞게)
        adjust_image_size(pic, img_width, img_height)

                    
            

        

     # 프레젠테이션을 저장할 경로 설정
    save_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint Files", "*.pptx")])

    if save_path:
        ppt.save(save_path)
        text.insert(tk.END, f"PowerPoint 프레젠테이션 저장 완료: {save_path}\n")
    text.insert(tk.END, "PowerPoint 프레젠테이션 생성 완료!\n")



    

def adjust_image_size(pic, target_width, target_height):
    pic_width = pic.width
    pic_height = pic.height
    aspect_ratio = pic_width / pic_height

    if pic_width > target_width:
        new_width = target_width
        new_height = int(new_width / aspect_ratio)
        pic.width = new_width
        pic.height = new_height

    if pic_height > target_height:
        new_height = target_height
        new_width = int(new_height * aspect_ratio)
        pic.width = new_width
        pic.height = new_height

def refresh():
    text.delete(1.0, tk.END)  # 텍스트 창 초기화

def create_ppt():
    # PowerPoint 파일을 생성하고, 파일 경로를 얻습니다.
    temp_dir = tempfile.gettempdir()
    ppt_file_path = os.path.join(temp_dir, "generated.pptx")
    
    # 여기에서 PowerPoint 파일 생성 작업을 수행하세요.
    # 예를 들어, python-pptx 라이브러리를 사용하여 파일을 생성할 수 있습니다.
    
    # 생성된 파일의 경로를 표시합니다.
    text.delete(1.0, tk.END)
    text.insert(tk.END, f"PowerPoint 파일 경로: {ppt_file_path}\n")

root = TkinterDnD.Tk()
root.title("PPTQuick vol.01")


image_path = resource_path("123.png")
image = tk.PhotoImage(file=image_path)

# 파일 선택 버튼 생성


label = tk.Label(root, image=image)
label.pack()

# 파일 선택 버튼 생성
select_button = tk.Button(root, text="환자 폴더 선택", command=open_folder_dialog, width=25, height=3)
select_button.pack(padx=20, pady=20)


text = tk.Text(root, wrap=tk.WORD, width=45, height=5)
text.pack(padx=10, pady=10)
text.insert(tk.END, "환자 폴더를 선택해주시고 pptx 저장 위치를 선택해주세요.\n")

# # "새로고침" 버튼 추가
# refresh_button = tk.Button(root, text="새로고침", command=refresh)
# refresh_button.pack()



# Label 위젯에 드래그 앤 드롭 기능 추가
label.drop_target_register(DND_FILES)
# label.dnd_bind('<<Drop>>', on_drop)
root.mainloop()